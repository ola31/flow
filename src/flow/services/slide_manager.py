"""SlideManager - PPTX 슬라이드를 이미지로 관리하는 서비스"""

from __future__ import annotations

import queue
import time
from collections import OrderedDict
from pathlib import Path

from pptx import Presentation
from PySide6.QtCore import QObject, QThread, Signal
from watchdog.events import FileSystemEventHandler
from watchdog.observers import Observer

from flow.services.slide_converter import (
    MarkdownSlideConverter,
    NoConverterAvailableError,
    SlideConverter,
    create_slide_converter,
)


class SlideLoadError(Exception):
    """PPTX 로드 실패 예외"""

    pass


class SlideUpdateHandler(FileSystemEventHandler):
    """파일 변경 이벤트 핸들러"""

    def __init__(self, target_path, callback):
        self.target_path = Path(target_path).resolve()
        self.callback = callback
        self.last_triggered = 0

    def on_modified(self, event):
        if event.is_directory:
            return

        if Path(event.src_path).resolve() != self.target_path:
            return

        now = time.time()
        if now - self.last_triggered > 0.1:
            self.callback()
            self.last_triggered = now


class PPTTask:
    """PPT 작업 단위 (큐에 담길 객체)"""

    LOAD_SINGLE = "LOAD_SINGLE"
    LOAD_METADATA = "LOAD_METADATA"

    def __init__(self, task_type, data):
        self.task_type = task_type
        self.data = data


class SlideWorker(QThread):
    """모든 PPT 작업을 순차적으로 처리하는 전용 백그라운드 스레드"""

    # 결과 전송용 시그널
    single_load_finished = Signal(int)
    metadata_finished = Signal(list)
    progress = Signal(int, int, str)
    status = Signal(str)
    error = Signal(str)

    def __init__(self, converter: SlideConverter):
        super().__init__()
        self._converter = converter
        self._task_queue = queue.Queue()
        self._is_running = True
        self._abort_requested = False

    def add_task(self, task: PPTTask):
        self._abort_requested = True
        while not self._task_queue.empty():
            try:
                self._task_queue.get_nowait()
            except queue.Empty:
                break
        self._abort_requested = False
        self._task_queue.put(task)

    def queue_task(self, task: PPTTask):
        """큐를 비우지 않고 뒤에 추가 (메타데이터 후 연속 변환 예약용).

        add_task는 진행 중 작업을 중단하고 큐를 비우므로, 여러 파일의
        변환을 순차 예약할 때는 이것을 써야 한다.
        """
        self._task_queue.put(task)

    def abort_current_task(self):
        self._abort_requested = True
        while not self._task_queue.empty():
            try:
                self._task_queue.get_nowait()
            except queue.Empty:
                break

    def stop(self):
        self._is_running = False
        self._abort_requested = True
        self.requestInterruption()
        # 유휴 워커는 큐 get(timeout=0.5)에 잠들어 있다 — 센티널로 즉시
        # 깨우지 않으면 GUI 스레드의 wait()가 최대 0.5초 블로킹된다
        # (홈 버튼마다 reset_worker가 이 경로를 밟음).
        self._task_queue.put(None)
        self.wait(1000)

    def run(self):
        while self._is_running:
            try:
                task = self._task_queue.get(timeout=0.5)
            except queue.Empty:
                if self.isInterruptionRequested():
                    break
                continue

            if task is None:  # stop()의 깨우기 센티널
                continue

            self._abort_requested = False
            try:
                if task.task_type == PPTTask.LOAD_SINGLE:
                    self._handle_single_load(task.data)
                elif task.task_type == PPTTask.LOAD_METADATA:
                    self._handle_metadata_load(task.data)
            except Exception as e:
                if not self._abort_requested and not self.isInterruptionRequested():
                    self.error.emit(str(e))
            finally:
                self._task_queue.task_done()

    def _count_slides(self, path: Path) -> int:
        """Counts slides using the converter for .md, python-pptx for .pptx.

        Only the MarkdownSlideConverter exposes get_slide_count; the .md
        branch is therefore only reached when this worker holds one.
        """
        if str(path).lower().endswith(".md"):
            converter = self._converter
            count = converter.get_slide_count(path)  # type: ignore[attr-defined]
            return int(count)
        prs = Presentation(str(path))
        return len(prs.slides)

    def _handle_single_load(self, path: Path):
        self.status.emit("PPT 파일 읽기 중...")
        try:
            slide_count = self._count_slides(path)
            engine_info = self._converter.get_engine_name()

            if slide_count > 0:
                for i in range(slide_count):
                    if self._abort_requested or self.isInterruptionRequested():
                        return
                    self._converter.convert_slide(
                        path, i, status_callback=self.status.emit
                    )
                    self.progress.emit(i + 1, slide_count, engine_info)

            if not self._abort_requested and not self.isInterruptionRequested():
                self.single_load_finished.emit(slide_count)
        except Exception as e:
            if not self._abort_requested:
                raise SlideLoadError(f"PPTX 로드 중 오류 발생: {e}")

    def _handle_metadata_load(self, song_data_list: list[tuple[str, Path]]):
        results = []
        for name, abs_p in song_data_list:
            if self._abort_requested or self.isInterruptionRequested():
                return
            count = 0
            try:
                count = self._count_slides(abs_p)
            except Exception:
                pass
            results.append((name, count))

        # 변환은 여기서 하지 않는다 — 카운트가 끝나야 오프셋/인덱스
        # globalize가 가능한데, 큰 PPT 변환을 기다리면 그동안 프로젝트의
        # 모든 핫스팟이 동작하지 않는다. 변환은 metadata_finished 후
        # SlideManager가 queue_task로 순차 예약한다.
        if not self._abort_requested and not self.isInterruptionRequested():
            self.metadata_finished.emit(results)


class SlideManager(QObject):
    """PPTX 파일을 로드하고 슬라이드 이미지를 관리함"""

    file_changed = Signal()
    load_started = Signal()
    load_finished = Signal(int)
    load_error = Signal(str)
    load_progress = Signal(int, int, str)
    load_status = Signal(str)

    songs_metadata_started = Signal()
    songs_metadata_finished = Signal(int)

    # PPT 변환 엔진(PowerPoint/LibreOffice)이 없을 때 PPT 조작 시도 시 발화.
    # MainWindow가 catch해서 설치 안내 다이얼로그를 띄운다.
    engine_missing = Signal()

    def __init__(self, converter: SlideConverter = None) -> None:
        super().__init__()
        self._pptx_path: Path | None = None
        self._slide_count: int = 0
        self._songs: list = []
        self._slide_offsets: dict[str, int] = {}
        self._total_slide_count: int = 0
        try:
            self._converter = converter or create_slide_converter()
        except NoConverterAvailableError:
            self._converter = None
        self._observer = None
        self._old_workers: list[SlideWorker] = []
        self._pending_reload_song = None
        # load_songs가 md/pptx 두 워커로 배치를 나눌 때 아직 안 끝난 배치 수.
        # 0이 될 때만 songs_metadata_finished를 발사한다 (두 번 발사되면
        # 핸들러의 인덱스 globalize가 중복 적용돼 매핑이 깨짐).
        self._pending_metadata_batches = 0
        # 백그라운드 변환/카운트 진행 여부 — peek_slide_image의 재예약 가드
        self._loading = False
        # 메타데이터 완료 후 순차 예약된 변환(LOAD_SINGLE) 잔여 수
        self._pending_conversions = 0
        # 이번 라운드에 워커 큐로 보낸 파일 — 같은 파일을 중복 예약하지 않기
        # 위한 것. _pending_conversions가 0이 되면 함께 비운다.
        self._queued_conversions: set[Path] = set()
        # 이번 load_songs에서 카운트한 곡 이름 (변환 예약 대상)
        self._counted_song_names: set[str] = set()
        # 변환 실패 후 자동 재예약 차단 플래그 — peek이 실패 파일을 계속
        # 재예약하면 에러 팝업이 무한 반복된다. 수동 로드만 재시도 허용.
        self._auto_retry_blocked = False
        # peek용 디코드 캐시 (키 → QImage). 핫스팟 클릭마다 같은 슬라이드
        # PNG를 재디코드하지 않게 한다. 키에 mtime이 있어 파일이 바뀌면
        # 자연히 미스가 나고, 상한으로 메모리를 제한한다.
        self._peek_lru: OrderedDict = OrderedDict()
        self._PEEK_LRU_MAX = 8
        # 미리보기/매핑 패널용 축소본 캐시 ((키, w, h) → QImage).
        # 클릭마다 풀해상도 스무스 스케일을 반복하지 않게 한다.
        self._thumb_lru: OrderedDict = OrderedDict()
        # 큰 프로젝트(수십~수백 장)도 전부 담기도록 넉넉히 — 장당 ~0.5MB
        self._THUMB_LRU_MAX = 192
        self.load_error.connect(self._on_load_error_clear_flag)

        if self._converter is not None:
            self._worker = SlideWorker(self._converter)
            self._connect_worker(self._worker)
            self._worker.start()
        else:
            self._worker = None

        # Markdown converter is always available — no external deps
        self._markdown_converter = MarkdownSlideConverter()
        self._markdown_worker = SlideWorker(self._markdown_converter)
        self._connect_worker(self._markdown_worker)
        self._markdown_worker.start()

        # 외부 PowerPoint 편집 중 파일 watcher 일시 중지 플래그
        self._watch_paused: bool = False

    def is_engine_available(self) -> bool:
        """PPT 변환 엔진이 사용 가능한지."""
        return self._converter is not None

    def stop_workers(self):
        for worker in (self._worker, self._markdown_worker):
            if worker is not None:
                worker.abort_current_task()
        # 큐를 통째로 비웠으니 예약 카운터도 함께 되돌린다. 안 그러면
        # 예약분에 대한 single_load_finished가 영영 오지 않아
        # _pending_conversions가 0으로 못 내려가고 _loading이 True로 굳는다
        # → 이후 모든 백그라운드 변환 예약이 조용히 무시되고, 라이브에서
        # 핫스팟을 골라도 슬라이드가 안 바뀌는 증상이 된다.
        self._pending_conversions = 0
        self._loading = False
        self._queued_conversions.clear()

    def clear_caches(self) -> None:
        """Clear cached converted slide images for all converter backends."""
        if self._converter is not None:
            self._converter.clear_cache()
        self._markdown_converter.clear_cache()
        self._peek_lru.clear()
        self._thumb_lru.clear()

    def is_watch_paused(self) -> bool:
        return self._watch_paused

    def pause_file_watching(self) -> None:
        """파일 watcher 일시 중지.

        외부 프로그램(예: PowerPoint)에서 PPT를 편집하는 동안 호출.
        편집 중 저장이 발생해도 자동 리로드를 시도하지 않아 파일 락 충돌
        및 PowerPoint 크래시를 예방한다.
        """
        self._watch_paused = True
        self.stop_watching()  # observer 자체 중지
        # 진행 중인 변환도 함께 중단해 COM/file handle 점유 해제
        self.stop_workers()

    def resume_file_watching(self) -> None:
        """파일 watcher 재개.

        사용자가 외부 편집을 마쳤을 때 호출. 보통 슬라이드 새로고침
        흐름과 함께 트리거된다.
        """
        self._watch_paused = False
        if self._pptx_path and self._pptx_path.exists():
            self.start_watching()

    def _connect_worker(self, worker: SlideWorker) -> None:
        worker.single_load_finished.connect(self._on_single_load_finished)
        worker.metadata_finished.connect(self._on_metadata_loaded)
        worker.progress.connect(self.load_progress.emit)
        worker.status.connect(self.load_status.emit)
        worker.error.connect(self.load_error.emit)

    def _disconnect_worker(self, worker: SlideWorker) -> None:
        worker.single_load_finished.disconnect(self._on_single_load_finished)
        worker.metadata_finished.disconnect(self._on_metadata_loaded)
        worker.progress.disconnect(self.load_progress.emit)
        worker.status.disconnect(self.load_status.emit)
        worker.error.disconnect(self.load_error.emit)

    def reset_worker(self):
        if self._converter is None:
            self._songs = []
            self._slide_offsets = {}
            self._total_slide_count = 0
            self._slide_count = 0
            return

        old = self._worker
        self._disconnect_worker(old)
        old.stop()

        self._old_workers = [w for w in self._old_workers if w.isRunning()]
        self._old_workers.append(old)

        self._songs = []
        self._slide_offsets = {}
        self._total_slide_count = 0
        self._slide_count = 0
        self._pending_metadata_batches = 0
        # 옛 워커를 버렸으니 그 큐에 남아 있던 예약도 없던 일이 된다 —
        # 카운터를 안 되돌리면 _loading이 True로 굳어 이후 변환 예약이
        # 전부 무시된다.
        self._pending_conversions = 0
        self._loading = False
        self._queued_conversions = set()

        self._worker = SlideWorker(self._converter)
        self._connect_worker(self._worker)
        self._worker.start()

    def _worker_for(self, path: Path) -> SlideWorker | None:
        """Return the worker matching the file's extension."""
        if str(path).lower().endswith(".md"):
            return self._markdown_worker
        return self._worker

    def load_pptx(self, path: str | Path):
        p = Path(path).resolve() if path and str(path).strip() else None
        if not p or not p.is_file():
            self._pptx_path = None
            self._slide_count = 0
            self.load_finished.emit(0)
            return

        worker = self._worker_for(p)
        if worker is None:
            self.engine_missing.emit()
            self.load_finished.emit(0)
            return

        self._pptx_path = p
        self._auto_retry_blocked = False
        # add_task는 큐를 비운다 — 예약 카운터도 이 한 건으로 맞춰 둬야
        # _loading이 True로 굳지 않는다.
        self._loading = True
        self._pending_conversions = 1
        self._queued_conversions = {p}
        self.load_started.emit()
        worker.add_task(PPTTask(PPTTask.LOAD_SINGLE, p))

    def _on_load_error_clear_flag(self, _msg: str) -> None:
        # 실패 시 로딩 플래그는 풀되, 자동 재예약은 차단한다 — 실패가
        # 반복 재시도(→ 에러 팝업 폭풍)로 이어지면 안 됨. 사용자 새로고침
        # (load_songs/reload_song/load_pptx)이 차단을 해제한다.
        self._loading = False
        self._pending_conversions = 0
        self._queued_conversions.clear()
        self._auto_retry_blocked = True

    def _on_single_load_finished(self, count: int):
        self._pending_conversions = max(0, self._pending_conversions - 1)
        self._loading = self._pending_conversions > 0
        if not self._loading:
            self._queued_conversions.clear()
        self._slide_count = count
        if self._pending_reload_song:
            self._pending_reload_song.set_slide_count(count)
            self._pending_reload_song = None
            self._recalculate_offsets()
        self.load_finished.emit(self.get_slide_count())

    def load_songs(self, songs: list, skip_counted: bool = False):
        """곡 목록 로드 (슬라이드 개수 카운트 + 오프셋 계산).

        skip_counted=True면 이미 카운트된 곡은 건너뛰고 새 곡만 카운트한다
        (곡 추가 시 전체 재로딩 방지). 파일 변경 반영이 필요한 전체 새로고침은
        기본값(False)으로 호출할 것.
        """
        self._songs = songs
        self._slide_offsets = {}
        self._total_slide_count = 0
        self._pending_metadata_batches = 0
        self._loading = False
        self._pending_conversions = 0
        self._queued_conversions = set()
        self._counted_song_names = set()
        self._auto_retry_blocked = False  # 사용자 주도 로드 = 재시도 허용

        song_data_list = []
        for s in songs:
            if skip_counted and s.get_slide_count() > 0:
                continue
            # markdown wins over pptx per Song.slide_source
            source = getattr(s, "slide_source", None)
            if source == "markdown":
                song_data_list.append((s.name, s.markdown_path))
            elif source == "pptx" or (source is None and s.has_slides):
                song_data_list.append((s.name, s.abs_slides_path))

        if not song_data_list:
            # 카운트할 곡이 없음 — 기존 카운트로 오프셋만 재계산.
            # songs_metadata_finished도 반드시 발생시켜야 함: 호출 측
            # (_on_songs_changed)이 localize 후 이 신호의 핸들러에서
            # 인덱스를 다시 globalize하기 때문.
            self._recalculate_offsets()
            self.songs_metadata_finished.emit(self._total_slide_count)
            self.load_finished.emit(self._total_slide_count)
            return

        # Split batch by file extension — markdown songs go to markdown worker,
        # pptx songs to pptx worker. Each worker only knows how to count its
        # own format.
        md_batch = [
            (n, p) for n, p in song_data_list
            if str(p).lower().endswith(".md")
        ]
        pptx_batch = [
            (n, p) for n, p in song_data_list
            if not str(p).lower().endswith(".md")
        ]

        if pptx_batch and self._worker is None:
            self.engine_missing.emit()
            self.load_finished.emit(0)
            return

        self._pending_metadata_batches = (1 if md_batch else 0) + (
            1 if pptx_batch else 0
        )
        self._loading = True
        self.songs_metadata_started.emit()
        if md_batch:
            self._markdown_worker.add_task(
                PPTTask(PPTTask.LOAD_METADATA, md_batch)
            )
        if pptx_batch:
            self._worker.add_task(
                PPTTask(PPTTask.LOAD_METADATA, pptx_batch)
            )

    def _on_metadata_loaded(self, results: list[tuple[str, int]]):
        if not self._songs:
            return

        for name, count in results:
            song = next((s for s in self._songs if s.name == name), None)
            if song:
                song.set_slide_count(count)

        self._counted_song_names.update(name for name, _ in results)

        # md/pptx 배치가 모두 끝났을 때만 완료 신호 (이중 globalize 방지)
        self._pending_metadata_batches = max(0, self._pending_metadata_batches - 1)
        if self._pending_metadata_batches > 0:
            return

        self._recalculate_offsets()

        # 이번에 카운트된 PPT 곡의 이미지 변환을 순차 예약 (백그라운드).
        # 카운트/오프셋은 이미 끝났으므로 UI는 즉시 동작하고, 썸네일은
        # 변환이 끝나는 대로 채워진다.
        self._pending_conversions = 0
        self._queued_conversions = set()
        if self._worker is not None:
            for s in self._songs:
                if s.name not in self._counted_song_names:
                    continue
                source = getattr(s, "slide_source", None)
                if source == "pptx" or (source is None and s.has_slides):
                    self._worker.queue_task(
                        PPTTask(PPTTask.LOAD_SINGLE, s.abs_slides_path)
                    )
                    self._queued_conversions.add(Path(s.abs_slides_path))
                    self._pending_conversions += 1
        self._loading = self._pending_conversions > 0

        self.songs_metadata_finished.emit(self._total_slide_count)
        self.load_finished.emit(self._total_slide_count)

    def get_slide_count(self) -> int:
        if self._total_slide_count > 0:
            return self._total_slide_count
        return self._slide_count

    def get_slide_image(self, index: int, status_callback=None):
        if self._total_slide_count > 0:
            try:
                song_name, local_index = self.global_to_local(index)
                return self.get_song_slide_image(
                    song_name, local_index, status_callback=status_callback
                )
            except Exception:
                return None

        # 단일 파일 모드 — _pptx_path 확장자에 따라 컨버터 선택
        if not self._pptx_path:
            return None
        if str(self._pptx_path).lower().endswith(".md"):
            return self._markdown_converter.convert_slide(
                self._pptx_path, index, status_callback=status_callback
            )
        if self._converter is None:
            return None
        return self._converter.convert_slide(
            self._pptx_path, index, status_callback=status_callback
        )

    def peek_slide_image(self, index: int, *, schedule: bool = True):
        """캐시된 슬라이드 이미지 즉시 반환 (변환 대기 없음, GUI 스레드용).

        캐시 미스면 None을 반환하고, schedule=True면 해당 파일의 변환을
        워커에 예약한다 — 완료되면 load_finished가 발사되므로 UI는 그때
        다시 조회해 채워 넣으면 된다.
        get_slide_image의 인라인 변환은 큰 PPT에서 UI를 얼리므로 GUI
        스레드에서는 반드시 이것을 쓸 것.

        Args:
            schedule: False면 캐시만 들여다보고 변환을 예약하지 않는다.
                짧은 주기로 완료를 기다리는 폴링에서 쓴다 — 폴링마다
                예약하면 PowerPoint 변환이 초당 몇 번씩 다시 돌아
                작은 창이 계속 떴다 사라진다.
        """
        lru_key = self.get_slide_cache_key(index)
        if lru_key is not None and lru_key in self._peek_lru:
            self._peek_lru.move_to_end(lru_key)
            return self._peek_lru[lru_key]

        if self._total_slide_count > 0:
            try:
                song_name, local_index = self.global_to_local(index)
            except Exception:
                return None
            song = next((s for s in self._songs if s.name == song_name), None)
            if song is None:
                return None
            source = getattr(song, "slide_source", None)
            if source == "markdown":
                # Qt 렌더라 인라인 안전 (컨버터가 자체 identity 캐시 보유)
                return self._markdown_converter.convert_slide(
                    song.markdown_path, local_index
                )
            if source == "pptx" or (source is None and song.has_slides):
                if self._converter is None:
                    return None
                img = self._converter.get_cached_slide(
                    song.abs_slides_path, local_index
                )
                if img is None:
                    if schedule:
                        self._ensure_background_conversion(song.abs_slides_path)
                else:
                    self._store_peek_lru(lru_key, img)
                return img
            return None

        # 단일 파일 모드
        if not self._pptx_path:
            return None
        if str(self._pptx_path).lower().endswith(".md"):
            return self._markdown_converter.convert_slide(self._pptx_path, index)
        if self._converter is None:
            return None
        img = self._converter.get_cached_slide(self._pptx_path, index)
        if img is None:
            if schedule:
                self._ensure_background_conversion(self._pptx_path)
        else:
            self._store_peek_lru(lru_key, img)
        return img

    def peek_thumbnail(self, index: int, max_w: int = 480, max_h: int = 270):
        """캐시된 슬라이드의 축소본을 즉시 반환 (GUI 스레드용, 비차단).

        PIP 미리보기·매핑 패널·팝오버처럼 작은 미리보기만 필요한 곳에서
        peek_slide_image 대신 사용 — 풀해상도 디코드/스케일을 반복하지
        않는다. 미변환이면 None (peek와 동일한 백그라운드 예약 동작).
        """
        from PySide6.QtCore import Qt

        base_key = self.get_slide_cache_key(index)
        tkey = (base_key, max_w, max_h) if base_key is not None else None
        if tkey is not None and tkey in self._thumb_lru:
            self._thumb_lru.move_to_end(tkey)
            return self._thumb_lru[tkey]

        img = self.peek_slide_image(index)
        if img is None or img.isNull():
            return None
        thumb = img.scaled(
            max_w,
            max_h,
            Qt.AspectRatioMode.KeepAspectRatio,
            Qt.TransformationMode.SmoothTransformation,
        )
        if tkey is not None:
            self._thumb_lru[tkey] = thumb
            self._thumb_lru.move_to_end(tkey)
            while len(self._thumb_lru) > self._THUMB_LRU_MAX:
                self._thumb_lru.popitem(last=False)
        return thumb

    def _store_peek_lru(self, key, img) -> None:
        if key is None:
            return
        self._peek_lru[key] = img
        self._peek_lru.move_to_end(key)
        while len(self._peek_lru) > self._PEEK_LRU_MAX:
            self._peek_lru.popitem(last=False)

    def get_slide_cache_key(self, index: int):
        """슬라이드 썸네일 캐시용 안정 키. 해석 불가면 None.

        (원본 파일 경로, mtime, 패치 mtime, 로컬 인덱스) — 파일이 안 바뀌면
        같은 키가 유지되므로, UI는 디코드 없이 캐시된 썸네일을 재사용할 수
        있다. 파일/패치가 바뀌면 mtime이 달라져 자동으로 무효화된다.
        """
        try:
            if self._total_slide_count > 0:
                song_name, local_index = self.global_to_local(index)
                song = next(
                    (s for s in self._songs if s.name == song_name), None
                )
                if song is None:
                    return None
                source = getattr(song, "slide_source", None)
                if source == "markdown":
                    src = song.markdown_path
                elif source == "pptx" or (source is None and song.has_slides):
                    src = song.abs_slides_path
                else:
                    return None
            else:
                if not self._pptx_path:
                    return None
                src = self._pptx_path
                local_index = index

            patches = src.parent / ".patches.json"
            patches_mtime = patches.stat().st_mtime if patches.exists() else 0.0
            return (str(src), src.stat().st_mtime, patches_mtime, local_index)
        except Exception:
            return None

    def register_appended_song(self, song) -> int:
        """라이브 중 셋리스트 끝에 추가된 곡을 송출 무중단으로 등록.

        슬라이드 카운트 + 오프셋 재계산만 수행하고 그 곡의 전역 오프셋을
        반환한다 (호출 측이 shift_indices로 매핑을 전역화). 완료 신호는
        발신하지 않는다 — songs_metadata_finished를 쏘면 핸들러의
        globalize가 이미 전역화된 기존 곡 인덱스를 이중 시프트한다.
        pptx면 이미지 변환을 백그라운드 큐 뒤에 예약해 둔다.
        """
        if not any(s is song for s in self._songs):
            # 빈 프로젝트로 열려 매니저 목록이 프로젝트와 비동기화된 경우
            self._songs.append(song)

        source = getattr(song, "slide_source", None)
        count = 0
        try:
            if source == "markdown":
                count = int(
                    self._markdown_converter.get_slide_count(song.markdown_path)
                )
            elif source == "pptx" or (source is None and song.has_slides):
                count = len(Presentation(str(song.abs_slides_path)).slides)
        except Exception:
            count = 0
        song.set_slide_count(count)
        self._recalculate_offsets()

        # pptx 이미지 변환은 백그라운드로 (큐를 비우지 않는 queue_task)
        if (
            count > 0
            and source != "markdown"
            and self._worker is not None
        ):
            self._worker.queue_task(
                PPTTask(PPTTask.LOAD_SINGLE, song.abs_slides_path)
            )
            self._queued_conversions.add(Path(song.abs_slides_path))
            self._pending_conversions += 1
            self._loading = True

        return self.get_song_offset(song.name)

    def _ensure_background_conversion(self, path) -> None:
        """해당 파일의 전체 변환을 워커에 예약 (중복 예약은 무시).

        load_started를 발신하지 않는다 — 조용한 백그라운드 워밍이므로
        로딩 오버레이(썸네일 영역 가림)를 띄우면 안 된다. 진행 표시는
        워커의 progress 신호가 패널 제목에 비차단으로 반영한다.

        다른 파일을 변환 중이더라도 요청을 버리지 않고 큐 뒤에 붙인다 —
        버리면 방금 필요해진 슬라이드가 영영 변환되지 않을 수 있다.
        """
        if self._worker is None or self._auto_retry_blocked:
            return
        path = Path(path)
        if path in self._queued_conversions:
            return
        if self._loading:
            self._queued_conversions.add(path)
            self._pending_conversions += 1
            self._worker.queue_task(PPTTask(PPTTask.LOAD_SINGLE, path))
            return
        self._loading = True
        self._pending_conversions = 1
        self._queued_conversions = {path}
        self._worker.add_task(PPTTask(PPTTask.LOAD_SINGLE, path))

    def get_song_slide_image(
        self, song_name: str, local_index: int, status_callback=None
    ):
        song = next((s for s in self._songs if s.name == song_name), None)
        if song is None:
            return None
        source = getattr(song, "slide_source", None)
        if source == "markdown":
            return self._markdown_converter.convert_slide(
                song.markdown_path, local_index, status_callback=status_callback
            )
        if source == "pptx" or (source is None and song.has_slides):
            if self._converter is None:
                return None
            return self._converter.convert_slide(
                song.abs_slides_path, local_index, status_callback=status_callback
            )
        return None

    def start_watching(self, path: str | Path = None):
        if path:
            self._pptx_path = Path(path)
        if not self._pptx_path or not self._pptx_path.parent.exists():
            return
        # pause 상태에선 다른 코드 경로(load_pptx 등)에서 호출돼도 watcher 비활성 유지
        if self._watch_paused:
            return

        resolved = self._pptx_path.resolve()
        # 같은 파일을 이미 감시 중이면 그대로 둔다 — 방향키 곡 전환마다
        # 옵저버 재시작(정지+조인 ~35ms)이 GUI를 막는 것 방지
        if (
            self._observer is not None
            and getattr(self, "_watched_path", None) == resolved
        ):
            self._pptx_path = resolved
            return

        self.stop_watching()
        self._pptx_path = resolved
        self._watched_path = resolved
        self._observer = Observer()
        handler = SlideUpdateHandler(self._pptx_path, self._on_watch_event)
        self._observer.schedule(handler, str(self._pptx_path.parent), recursive=False)
        self._observer.start()

    def _on_watch_event(self) -> None:
        """Watcher callback. Invalidates markdown cache before emitting.

        For .md files, the in-memory render cache must be cleared before any
        listener responds to ``file_changed`` so a re-render reads fresh
        content from disk.
        """
        if self._pptx_path is not None and str(self._pptx_path).lower().endswith(
            ".md"
        ):
            self._markdown_converter.invalidate_cache(self._pptx_path)
        self.file_changed.emit()

    def stop_watching(self):
        if self._observer:
            # join으로 GUI를 막지 않는다 — 곡 전환마다 ~35ms씩 걸리던 비용.
            # 정지 요청만 하고 스레드는 뒤에서 끝나게 둔다 (shutdown에서 join).
            self._observer.stop()
            self._dying_observers = [
                o for o in getattr(self, "_dying_observers", [])
                if o.is_alive()
            ]
            self._dying_observers.append(self._observer)
            self._observer = None
        self._watched_path = None

    def shutdown(self):
        self.stop_watching()
        # 앱 종료 시엔 비동기로 정지 요청한 옵저버들도 확실히 정리
        for obs in getattr(self, "_dying_observers", []):
            obs.join(timeout=1)
        self._dying_observers = []
        for worker in (self._worker, self._markdown_worker):
            if worker is not None:
                worker.stop()
        for w in self._old_workers:
            if w.isRunning():
                w.stop()
        self._old_workers.clear()

    def invalidate_markdown_cache(self, md_path: Path) -> None:
        """Public hook to drop the markdown render cache for one song."""
        self._markdown_converter.invalidate_cache(md_path)

    def global_to_local(self, global_index: int) -> tuple[str, int]:
        for song in self._songs:
            offset = self._slide_offsets.get(song.name, 0)
            count = song.get_slide_count()
            if offset <= global_index < offset + count:
                return (song.name, global_index - offset)
        raise ValueError(f"Invalid index: {global_index}")

    def local_to_global(self, song_name: str, local_index: int) -> int:
        offset = self._slide_offsets.get(song_name)
        if offset is None:
            raise ValueError(f"Song not found: {song_name}")
        return offset + local_index

    def get_song_offset(self, song_name: str) -> int:
        return self._slide_offsets.get(song_name, 0)

    def _recalculate_offsets(self) -> None:
        offset = 0
        for song in self._songs:
            source = getattr(song, "slide_source", None)
            has_source = (
                source in ("markdown", "pptx")
                or (source is None and song.has_slides)
            )
            if has_source:
                self._slide_offsets[song.name] = offset
                offset += song.get_slide_count()
        self._total_slide_count = offset

    def reload_song(self, song):
        # markdown 우선, 없으면 PPT, 둘 다 없으면 0장 처리
        source = getattr(song, "slide_source", None)
        if source == "markdown":
            target_path = song.markdown_path
        elif source == "pptx" or (source is None and song.has_slides):
            target_path = song.abs_slides_path
        else:
            target_path = None

        if target_path is None:
            if self._songs:
                song.set_slide_count(0)
                self._recalculate_offsets()
                self.load_finished.emit(self._total_slide_count)
            return

        worker = self._worker_for(target_path)
        if worker is None:
            self.engine_missing.emit()
            return

        # markdown converter 캐시는 source 변경 가능성에 대비해 invalidate
        if source == "markdown":
            self._markdown_converter.invalidate_cache(target_path)

        if self._songs:
            self._pending_reload_song = song
        self._auto_retry_blocked = False  # 사용자 새로고침 = 재시도 허용
        # add_task는 큐를 비운다 — 예약 카운터도 이 한 건으로 맞춰야
        # 사라진 예약분 때문에 _loading이 True로 굳지 않는다.
        self._loading = True
        self._pending_conversions = 1
        self._queued_conversions = {Path(target_path)}
        self.load_started.emit()
        worker.add_task(PPTTask(PPTTask.LOAD_SINGLE, target_path))

    def reload_all_songs(self):
        if not self._songs:
            return
        # 곡 중에 PPT가 하나라도 있으면 PPT 엔진 필요
        has_pptx = any(
            getattr(s, "slide_source", None) == "pptx" or
            (getattr(s, "slide_source", None) is None and s.has_slides)
            for s in self._songs
        )
        if has_pptx and self._converter is None:
            self.engine_missing.emit()
            return
        # 캐시는 지우지 않는다 — 변환 캐시 키에 파일 mtime이 포함돼 있어
        # 변경된 파일은 자동으로 재변환되고(키가 달라짐), 안 바뀐 곡은
        # 캐시를 그대로 쓴다. 전부 지우면 전체 재변환으로 새로고침이
        # 수십 초씩 걸린다. (해상도 변경 시에는 별도 경로에서 clear_caches 호출)
        self.load_songs(self._songs)
