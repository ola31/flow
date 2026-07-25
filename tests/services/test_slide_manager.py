import time

import pytest
from unittest.mock import MagicMock, patch
from pathlib import Path

from flow.services.slide_manager import SlideManager, SlideWorker, PPTTask


@pytest.fixture
def mock_converter():
    converter = MagicMock()
    converter.convert_slide.return_value = None
    converter.get_engine_name.return_value = "Mock"
    return converter


@pytest.fixture
def manager(mock_converter):
    mgr = SlideManager(converter=mock_converter)
    yield mgr
    # Tear down everything (worker + markdown worker + file watcher) so that
    # an asserting test never leaks a QThread/Observer → SIGABRT during exit.
    mgr.shutdown()


class TestSlideManager:
    def test_initial_slide_count_is_zero(self, manager):
        assert manager.get_slide_count() == 0

    def test_get_slide_image_returns_qimage(self, manager, mock_converter):
        from PySide6.QtGui import QImage

        mock_image = QImage(100, 100, QImage.Format.Format_RGB32)
        mock_converter.convert_slide.return_value = mock_image

        manager._pptx_path = Path("/fake/test.pptx")
        image = manager.get_slide_image(0)

        assert isinstance(image, QImage)
        assert image.width() == 100

    def test_reset_worker_clears_state(self, manager):
        manager._songs = ["fake_song"]
        manager._slide_offsets = {"test": 5}
        manager._total_slide_count = 10

        manager.reset_worker()

        assert manager._songs == []
        assert manager._slide_offsets == {}
        assert manager._total_slide_count == 0
        assert manager._slide_count == 0

    def test_global_to_local_raises_on_invalid_index(self, manager):
        with pytest.raises(ValueError, match="Invalid index"):
            manager.global_to_local(999)

    def test_local_to_global_raises_on_unknown_song(self, manager):
        with pytest.raises(ValueError, match="Song not found"):
            manager.local_to_global("nonexistent", 0)

    def test_file_watcher_notifies_on_change(self, tmp_path, manager, qtbot):
        pptx_file = tmp_path / "test.pptx"
        pptx_file.write_text("initial content")

        manager.start_watching(pptx_file)
        # Give the watchdog Observer thread a moment to register the inotify
        # watch before we modify the file. Without this the modification can
        # race ahead of the watcher and the event is lost on slower runners.
        import time
        time.sleep(0.5)

        # qtbot.waitSignal pumps the Qt event loop while waiting, which is
        # required because file_changed is emitted from the watchdog thread
        # and arrives via QueuedConnection.
        with qtbot.waitSignal(manager.file_changed, timeout=4000):
            pptx_file.write_text("updated content")

        manager.stop_watching()


class TestSlideWorker:
    def test_abort_clears_queue(self, mock_converter):
        worker = SlideWorker(mock_converter)
        worker._task_queue.put(PPTTask(PPTTask.LOAD_SINGLE, Path("/fake")))
        worker._task_queue.put(PPTTask(PPTTask.LOAD_SINGLE, Path("/fake2")))

        worker.abort_current_task()

        assert worker._task_queue.empty()
        assert worker._abort_requested is True

    def test_add_task_resets_abort_flag(self, mock_converter):
        worker = SlideWorker(mock_converter)
        worker._abort_requested = True

        worker.add_task(PPTTask(PPTTask.LOAD_SINGLE, Path("/fake")))

        assert worker._abort_requested is False
        assert not worker._task_queue.empty()

    def test_stop_returns_quickly_when_idle(self, mock_converter):
        """stop()은 GUI 스레드에서 불린다 — 유휴 워커가 큐 get(timeout=0.5)에
        잠들어 있어도 즉시 깨워서 반환해야 한다. 안 그러면 홈 버튼 클릭마다
        (reset_worker 경로) UI가 최대 0.5초 멈춘다."""
        worker = SlideWorker(mock_converter)
        worker.start()
        time.sleep(0.1)  # 워커가 큐 대기에 들어간 시점

        t0 = time.time()
        worker.stop()
        elapsed = time.time() - t0

        assert elapsed < 0.2, f"stop()이 {elapsed:.2f}s 블로킹 — 즉시 깨워야 함"
        assert not worker.isRunning()


class TestFileWatcherPause:
    """외부 PPT 편집 시나리오를 위한 watcher 일시 중지/재개."""

    def test_initial_state_not_paused(self, manager):
        assert manager.is_watch_paused() is False

    def test_pause_sets_flag_and_stops_observer(self, manager, tmp_path):
        # observer가 동작 중이라고 가정 — pptx_path 설정 후 시작
        pptx = tmp_path / "x.pptx"
        pptx.write_bytes(b"fake")
        manager.start_watching(pptx)
        assert manager._observer is not None

        manager.pause_file_watching()
        assert manager.is_watch_paused() is True
        assert manager._observer is None  # stop_watching이 None으로 만듦

    def test_paused_state_blocks_start_watching(self, manager, tmp_path):
        pptx = tmp_path / "x.pptx"
        pptx.write_bytes(b"fake")
        manager.pause_file_watching()  # pause 먼저

        # 다른 코드 경로(load_pptx 등)에서 start_watching이 호출돼도 무시
        manager.start_watching(pptx)
        assert manager._observer is None

    def test_resume_restarts_watcher(self, manager, tmp_path):
        pptx = tmp_path / "x.pptx"
        pptx.write_bytes(b"fake")
        manager._pptx_path = pptx.resolve()

        manager.pause_file_watching()
        assert manager.is_watch_paused() is True

        manager.resume_file_watching()
        assert manager.is_watch_paused() is False
        assert manager._observer is not None

    def test_resume_noop_if_pptx_missing(self, manager, tmp_path):
        manager._pptx_path = tmp_path / "no.pptx"  # 존재하지 않음

        manager.pause_file_watching()
        manager.resume_file_watching()

        assert manager.is_watch_paused() is False
        assert manager._observer is None  # 파일 없으면 시작 안 함


class TestIncrementalLoadSongs:
    """곡 추가 시 전체 재카운트 방지 — skip_counted 증분 로드"""

    def _make_song(self, tmp_path, name, count):
        from flow.domain.song import Song

        song_dir = tmp_path / "songs" / name
        song_dir.mkdir(parents=True, exist_ok=True)
        (song_dir / "slides.pptx").write_bytes(b"PK fake")
        song = Song(name=name, folder=Path("songs") / name, project_dir=tmp_path)
        if count:
            song.set_slide_count(count)
        return song

    def test_skip_counted_only_queues_new_songs(self, manager, tmp_path):
        song_a = self._make_song(tmp_path, "song_a", 3)
        song_b = self._make_song(tmp_path, "song_b", 2)
        new_song = self._make_song(tmp_path, "song_new", 0)

        tasks = []
        manager._worker.add_task = tasks.append

        manager.load_songs([song_a, song_b, new_song], skip_counted=True)

        assert len(tasks) == 1
        batch = tasks[0].data
        assert [name for name, _ in batch] == ["song_new"]

    def test_skip_counted_all_counted_finishes_without_worker(
        self, manager, tmp_path, qtbot
    ):
        song_a = self._make_song(tmp_path, "song_a", 3)
        song_b = self._make_song(tmp_path, "song_b", 2)

        tasks = []
        manager._worker.add_task = tasks.append
        finished = []
        manager.load_finished.connect(finished.append)
        meta_finished = []
        manager.songs_metadata_finished.connect(meta_finished.append)

        manager.load_songs([song_a, song_b], skip_counted=True)

        assert tasks == []
        assert finished == [5]
        # 인덱스 글로벌화(_globalize_project_indices)는 songs_metadata_finished
        # 핸들러에서 일어나므로 이 신호도 반드시 발생해야 한다.
        assert meta_finished == [5]
        assert manager.get_song_offset("song_a") == 0
        assert manager.get_song_offset("song_b") == 3

    def test_default_recounts_everything(self, manager, tmp_path):
        song_a = self._make_song(tmp_path, "song_a", 3)
        new_song = self._make_song(tmp_path, "song_new", 0)

        tasks = []
        manager._worker.add_task = tasks.append

        manager.load_songs([song_a, new_song])

        assert len(tasks) == 1
        assert [name for name, _ in tasks[0].data] == ["song_a", "song_new"]


class TestReloadAllKeepsCaches:
    """전체 새로고침이 변환 캐시를 통째로 지우지 않아야 한다.

    캐시 키에 파일 mtime이 포함되므로 변경된 파일은 자동으로 재변환된다.
    캐시를 전부 지우면 안 바뀐 곡까지 전 슬라이드 재변환이 일어나
    새로고침이 극단적으로 느려진다.
    """

    def test_reload_all_does_not_clear_converter_caches(
        self, manager, mock_converter, tmp_path
    ):
        from pptx import Presentation

        from flow.domain.song import Song

        song_dir = tmp_path / "songs" / "song_a"
        song_dir.mkdir(parents=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(song_dir / "slides.pptx")
        song = Song(name="song_a", folder=Path("songs/song_a"), project_dir=tmp_path)
        song.set_slide_count(1)
        manager._songs = [song]

        md_cache_cleared = []
        manager._markdown_converter.clear_cache = lambda: md_cache_cleared.append(True)

        manager.reload_all_songs()

        mock_converter.clear_cache.assert_not_called()
        assert md_cache_cleared == []


class TestMixedFormatMetadataLoad:
    """마크다운+PPT 혼합 프로젝트 메타데이터 로드.

    배치가 두 워커로 갈라져도 songs_metadata_finished는 전체 완료 시
    정확히 한 번만 발생해야 한다 — 이 신호의 핸들러가 인덱스를
    globalize하므로 두 번 발생하면 매핑 인덱스가 이중 시프트로 깨진다.
    """

    def _make_md_song(self, tmp_path, name, n_slides):
        from flow.domain.song import Song

        song_dir = tmp_path / "songs" / name
        song_dir.mkdir(parents=True, exist_ok=True)
        body = "\n\n".join(f"lyric line {i}" for i in range(n_slides))
        (song_dir / "slides.md").write_text(f"# {name}\n\n{body}\n", encoding="utf-8")
        return Song(name=name, folder=Path("songs") / name, project_dir=tmp_path)

    def _make_pptx_song(self, tmp_path, name, n_slides):
        from pptx import Presentation

        from flow.domain.song import Song

        song_dir = tmp_path / "songs" / name
        song_dir.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        for _ in range(n_slides):
            prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(song_dir / "slides.pptx")
        return Song(name=name, folder=Path("songs") / name, project_dir=tmp_path)

    def test_metadata_finished_emitted_once_with_full_total(
        self, manager, tmp_path, qtbot
    ):
        md_song = self._make_md_song(tmp_path, "song_md", 2)
        pptx_song = self._make_pptx_song(tmp_path, "song_ppt", 3)

        meta = []
        manager.songs_metadata_finished.connect(meta.append)

        manager.load_songs([md_song, pptx_song])

        qtbot.waitUntil(lambda: len(meta) >= 1, timeout=15000)
        qtbot.wait(1000)  # 늦게 끝나는 배치가 신호를 한 번 더 쏘는지 관찰

        assert meta == [5]
        assert manager.get_song_offset("song_md") == 0
        assert manager.get_song_offset("song_ppt") == 2


class TestPeekSlideImage:
    """GUI 스레드용 비차단 슬라이드 조회.

    peek은 절대 인라인 변환을 하지 않는다 — 캐시 미스면 None을 반환하고
    (유휴 상태일 때만) 백그라운드 변환을 예약한다. 인라인 변환은 61장짜리
    PPT에서 UI를 수십 초 얼린다 ("Flow is not responding").
    """

    def _make_pptx_song(self, tmp_path, name="song_p"):
        from pptx import Presentation

        from flow.domain.song import Song

        d = tmp_path / "songs" / name
        d.mkdir(parents=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(d / "slides.pptx")
        song = Song(name=name, folder=Path("songs") / name, project_dir=tmp_path)
        song.set_slide_count(1)
        return song

    def test_peek_never_converts_inline(self, manager, mock_converter, tmp_path):
        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = None
        manager._loading = True  # 백그라운드 로드 중 → 재예약도 하지 않음

        img = manager.peek_slide_image(0)

        assert img is None
        mock_converter.convert_slide.assert_not_called()

    def test_peek_returns_cached_image(self, manager, mock_converter, tmp_path):
        from PySide6.QtGui import QImage

        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        cached = QImage(4, 4, QImage.Format.Format_RGB32)
        mock_converter.get_cached_slide.return_value = cached

        img = manager.peek_slide_image(0)

        assert img is cached
        mock_converter.convert_slide.assert_not_called()

    def test_peek_miss_schedules_background_conversion_when_idle(
        self, manager, mock_converter, tmp_path
    ):
        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = None
        tasks = []
        manager._worker.add_task = tasks.append

        manager.peek_slide_image(0)

        assert len(tasks) == 1  # LOAD_SINGLE 예약
        # 로딩 중 재호출은 중복 예약하지 않음
        manager.peek_slide_image(0)
        assert len(tasks) == 1

    def test_peek_miss_does_not_interrupt_active_load(
        self, manager, mock_converter, tmp_path
    ):
        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = None
        manager._loading = True
        tasks = []
        manager._worker.add_task = tasks.append

        manager.peek_slide_image(0)

        assert tasks == []


class TestSlideCacheKey:
    """썸네일 캐시용 안정 키 — 파일이 안 바뀌면 같고, 바뀌면 달라져야 한다."""

    def _make_pptx_song(self, tmp_path, name="song_k"):
        from pptx import Presentation

        from flow.domain.song import Song

        d = tmp_path / "songs" / name
        d.mkdir(parents=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(d / "slides.pptx")
        song = Song(name=name, folder=Path("songs") / name, project_dir=tmp_path)
        song.set_slide_count(2)
        return song

    def test_key_stable_for_unchanged_file(self, manager, tmp_path):
        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()

        assert manager.get_slide_cache_key(0) == manager.get_slide_cache_key(0)
        assert manager.get_slide_cache_key(0) != manager.get_slide_cache_key(1)

    def test_key_changes_when_file_modified(self, manager, tmp_path):
        import os

        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()

        before = manager.get_slide_cache_key(0)
        os.utime(song.abs_slides_path, (1e9, 1e9))
        after = manager.get_slide_cache_key(0)

        assert before != after

    def test_key_none_for_invalid_index(self, manager, tmp_path):
        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()

        assert manager.get_slide_cache_key(99) is None


class TestMetadataCountsBeforeConverting:
    """메타데이터 로드는 카운트만 빠르게 끝내고, 변환은 뒤로 미룬다.

    카운트+변환이 한 작업이면 61장 PPT 변환이 끝날 때까지 오프셋/인덱스
    globalize가 지연돼, 그동안 프로젝트의 모든 핫스팟이 동작하지 않는다.
    """

    def _pptx(self, tmp_path, name, slides=2):
        from pptx import Presentation

        d = tmp_path / "songs" / name
        d.mkdir(parents=True)
        prs = Presentation()
        for _ in range(slides):
            prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(d / "slides.pptx")
        return d / "slides.pptx"

    def test_metadata_task_does_not_convert(self, mock_converter, tmp_path, qtbot):
        worker = SlideWorker(mock_converter)
        p = self._pptx(tmp_path, "song_a", 3)
        results = []
        worker.metadata_finished.connect(results.append)

        worker._handle_metadata_load([("song_a", p)])

        assert results == [[("song_a", 3)]]
        mock_converter.convert_slide.assert_not_called()

    def test_conversion_enqueued_after_final_metadata(
        self, manager, mock_converter, tmp_path
    ):
        from flow.domain.song import Song

        p = self._pptx(tmp_path, "song_a", 2)
        song = Song(name="song_a", folder=Path("songs/song_a"), project_dir=tmp_path)
        manager._songs = [song]
        manager._pending_metadata_batches = 1
        manager._loading = True

        queued = []
        manager._worker.queue_task = queued.append

        manager._on_metadata_loaded([("song_a", 2)])

        assert len(queued) == 1 and queued[0].data == p
        # 변환이 남아 있는 동안 _loading 유지 (peek의 add_task가 큐를 비우는
        # 사고 방지)
        assert manager._loading is True

        manager._on_single_load_finished(2)
        assert manager._loading is False

    def test_worker_queue_task_does_not_clear_queue(self, mock_converter):
        worker = SlideWorker(mock_converter)
        worker._task_queue.put(PPTTask(PPTTask.LOAD_SINGLE, Path("/fake1")))

        worker.queue_task(PPTTask(PPTTask.LOAD_SINGLE, Path("/fake2")))

        assert worker._task_queue.qsize() == 2


class TestPeekImageLRU:
    """peek이 같은 슬라이드를 반복 조회할 때 PNG 재디코드를 피한다.

    핫스팟 클릭마다 PIP+매핑 패널이 같은 슬라이드를 다시 요청하는데,
    매번 풀해상도 PNG를 디코드하면 클릭이 수백 ms씩 걸린다.
    """

    def _song(self, tmp_path):
        from pptx import Presentation

        from flow.domain.song import Song

        d = tmp_path / "songs" / "song_l"
        d.mkdir(parents=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(d / "slides.pptx")
        song = Song(name="song_l", folder=Path("songs/song_l"), project_dir=tmp_path)
        song.set_slide_count(1)
        return song

    def test_repeated_peek_hits_lru(self, manager, mock_converter, tmp_path):
        from PySide6.QtGui import QImage

        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        img = QImage(4, 4, QImage.Format.Format_RGB32)
        mock_converter.get_cached_slide.return_value = img

        first = manager.peek_slide_image(0)
        second = manager.peek_slide_image(0)

        assert first is img and second is img
        assert mock_converter.get_cached_slide.call_count == 1

    def test_lru_invalidated_by_file_change(self, manager, mock_converter, tmp_path):
        import os

        from PySide6.QtGui import QImage

        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = QImage(
            4, 4, QImage.Format.Format_RGB32
        )

        manager.peek_slide_image(0)
        os.utime(song.abs_slides_path, (2e9, 2e9))  # mtime 변경 → 키 변경
        manager.peek_slide_image(0)

        assert mock_converter.get_cached_slide.call_count == 2

    def test_clear_caches_drops_lru(self, manager, mock_converter, tmp_path):
        from PySide6.QtGui import QImage

        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = QImage(
            4, 4, QImage.Format.Format_RGB32
        )

        manager.peek_slide_image(0)
        manager.clear_caches()
        manager.peek_slide_image(0)

        assert mock_converter.get_cached_slide.call_count == 2


class TestPeekThumbnail:
    """미리보기/매핑 패널용 축소본 조회 — 반복 클릭 시 스케일 재연산 방지."""

    def _song(self, tmp_path):
        from pptx import Presentation

        from flow.domain.song import Song

        d = tmp_path / "songs" / "song_t"
        d.mkdir(parents=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(d / "slides.pptx")
        song = Song(name="song_t", folder=Path("songs/song_t"), project_dir=tmp_path)
        song.set_slide_count(1)
        return song

    def test_thumbnail_scaled_within_bounds(self, manager, mock_converter, tmp_path):
        from PySide6.QtGui import QImage

        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = QImage(
            1920, 1080, QImage.Format.Format_RGB32
        )

        thumb = manager.peek_thumbnail(0, 480, 270)

        assert thumb is not None
        assert thumb.width() <= 480 and thumb.height() <= 270

    def test_repeat_thumbnail_skips_rescale_and_decode(
        self, manager, mock_converter, tmp_path
    ):
        from PySide6.QtGui import QImage

        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = QImage(
            1920, 1080, QImage.Format.Format_RGB32
        )

        first = manager.peek_thumbnail(0, 480, 270)
        peeks = []
        manager.peek_slide_image = lambda i: peeks.append(i)  # 재호출 감지
        second = manager.peek_thumbnail(0, 480, 270)

        assert second is first  # 캐시된 동일 객체
        assert peeks == []

    def test_thumbnail_none_when_unconverted(self, manager, mock_converter, tmp_path):
        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = None
        manager._loading = True

        assert manager.peek_thumbnail(0, 480, 270) is None


class TestConversionFailureBackoff:
    """변환 실패 시 자동 재예약 금지 — 실패 파일을 peek이 계속 재예약하면
    에러 팝업이 무한 반복돼 앱이 죽는다 (Windows에서 실제 발생).
    수동 로드(load_songs/reload_song)만 재시도를 허용한다."""

    def _song(self, tmp_path):
        from pptx import Presentation

        from flow.domain.song import Song

        d = tmp_path / "songs" / "song_f"
        d.mkdir(parents=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(d / "slides.pptx")
        song = Song(name="song_f", folder=Path("songs/song_f"), project_dir=tmp_path)
        song.set_slide_count(1)
        return song

    def test_error_blocks_auto_reschedule(self, manager, mock_converter, tmp_path):
        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = None
        tasks = []
        manager._worker.add_task = tasks.append

        manager.peek_slide_image(0)   # 1차 예약
        assert len(tasks) == 1
        manager.load_error.emit("변환 실패")  # 실패 → _loading 해제 + 재시도 차단

        manager.peek_slide_image(0)   # 자동 재예약되면 안 됨
        manager.peek_slide_image(0)
        assert len(tasks) == 1

    def test_manual_reload_clears_block(self, manager, mock_converter, tmp_path):
        song = self._song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = None
        tasks = []
        manager._worker.add_task = tasks.append

        manager.peek_slide_image(0)
        manager.load_error.emit("변환 실패")
        manager.reload_song(song)     # 사용자 새로고침 = 명시적 재시도

        assert len(tasks) == 2        # 재시도 허용
        # 재시도도 실패하면 다시 차단
        manager.load_error.emit("변환 실패")
        manager.peek_slide_image(0)
        assert len(tasks) == 2


class TestRegisterAppendedSong:
    """라이브 중 셋리스트 끝에 추가된 곡의 무중단 등록.

    카운트 + 오프셋 재계산만 수행 (기존 곡 오프셋 불변, 완료 신호 없음 —
    신호를 쏘면 핸들러의 globalize가 기존 곡 인덱스를 이중 시프트한다).
    """

    def _md_song(self, tmp_path, name, n):
        from flow.domain.song import Song

        d = tmp_path / "songs" / name
        d.mkdir(parents=True)
        body = "\n\n".join(f"line {i}" for i in range(n))
        (d / "slides.md").write_text(f"# {name}\n\n{body}\n", encoding="utf-8")
        return Song(name=name, folder=Path("songs") / name, project_dir=tmp_path)

    def _pptx_song(self, tmp_path, name, n):
        from pptx import Presentation

        from flow.domain.song import Song

        d = tmp_path / "songs" / name
        d.mkdir(parents=True)
        prs = Presentation()
        for _ in range(n):
            prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(d / "slides.pptx")
        return Song(name=name, folder=Path("songs") / name, project_dir=tmp_path)

    def test_appended_md_song_gets_count_and_offset(self, qapp, manager, tmp_path):
        existing = self._md_song(tmp_path, "song_a", 3)
        existing.set_slide_count(3)
        songs = [existing]
        manager._songs = songs
        manager._recalculate_offsets()

        new_song = self._md_song(tmp_path, "song_new", 2)
        songs.append(new_song)  # 라이브 중 append (aliased list)

        offset = manager.register_appended_song(new_song)

        assert new_song.get_slide_count() == 2
        assert offset == 3
        assert manager.get_song_offset("song_a") == 0  # 기존 곡 불변
        assert manager.get_slide_count() == 5

    def test_appended_pptx_song_counted(self, qapp, manager, tmp_path):
        existing = self._md_song(tmp_path, "song_a", 3)
        existing.set_slide_count(3)
        songs = [existing]
        manager._songs = songs
        manager._recalculate_offsets()

        new_song = self._pptx_song(tmp_path, "song_ppt", 4)
        songs.append(new_song)

        offset = manager.register_appended_song(new_song)

        assert new_song.get_slide_count() == 4
        assert offset == 3

    def test_no_completion_signals_emitted(self, qapp, manager, tmp_path):
        existing = self._md_song(tmp_path, "song_a", 3)
        existing.set_slide_count(3)
        songs = [existing]
        manager._songs = songs
        manager._recalculate_offsets()
        new_song = self._md_song(tmp_path, "song_new", 2)
        songs.append(new_song)

        fired = []
        manager.songs_metadata_finished.connect(lambda n: fired.append(n))
        manager.load_finished.connect(lambda n: fired.append(n))

        manager.register_appended_song(new_song)

        assert fired == []  # globalize 핸들러가 돌면 안 됨

    def test_untracked_song_gets_adopted(self, qapp, manager, tmp_path):
        """매니저 목록과 비동기화된 경우(빈 프로젝트로 연 라이브)도 등록."""
        new_song = self._md_song(tmp_path, "song_solo", 2)
        # manager._songs는 빈 별도 리스트인 상황

        offset = manager.register_appended_song(new_song)

        assert offset == 0
        assert manager.get_slide_count() == 2


class TestWatcherSamePathSkip:
    def test_same_path_does_not_restart_observer(self, manager, tmp_path):
        """방향키로 같은 곡 시트를 오갈 때마다 감시자를 재시작(조인 ~35ms)
        하면 전환이 느려진다 — 같은 파일이면 그대로 둔다."""
        pptx = tmp_path / "slides.pptx"
        pptx.touch()

        manager.start_watching(str(pptx))
        first = manager._observer
        assert first is not None

        manager.start_watching(str(pptx))

        assert manager._observer is first  # 재시작 없음
        manager.stop_watching()

    def test_different_path_restarts(self, manager, tmp_path):
        a = tmp_path / "a" / "slides.pptx"
        b = tmp_path / "b" / "slides.pptx"
        a.parent.mkdir(); b.parent.mkdir()
        a.touch(); b.touch()

        manager.start_watching(str(a))
        first = manager._observer
        manager.start_watching(str(b))

        assert manager._observer is not first
        manager.stop_watching()


class TestConversionQueueNeverStalls:
    """예약 카운터가 굳으면 백그라운드 변환이 조용히 멈춘다.

    _loading이 True로 남으면 이후 peek 미스가 변환을 예약하지 못하고,
    라이브에서 핫스팟을 눌러도 슬라이드가 안 바뀌는 증상이 된다.
    """

    def _make_pptx_song(self, tmp_path, name="song_q"):
        from pptx import Presentation

        from flow.domain.song import Song

        d = tmp_path / "songs" / name
        d.mkdir(parents=True)
        Presentation().save(str(d / "slides.pptx"))
        song = Song(name=name, folder=d, project_dir=tmp_path)
        song.set_slide_count(3)
        return song

    def test_stop_workers_resets_pending_counter(self, manager):
        manager._pending_conversions = 4
        manager._loading = True
        manager._queued_conversions = {Path("/a.pptx")}

        manager.stop_workers()

        assert manager._pending_conversions == 0
        assert manager._loading is False
        assert manager._queued_conversions == set()

    def test_reset_worker_resets_pending_counter(self, manager):
        manager._pending_conversions = 2
        manager._loading = True

        manager.reset_worker()

        assert manager._pending_conversions == 0
        assert manager._loading is False

    def test_busy_worker_queues_instead_of_dropping(
        self, manager, mock_converter, tmp_path
    ):
        song = self._make_pptx_song(tmp_path)
        manager._songs = [song]
        manager._recalculate_offsets()
        mock_converter.get_cached_slide.return_value = None
        manager._loading = True
        manager._pending_conversions = 1
        manager._queued_conversions = set()
        queued = []
        manager._worker.queue_task = queued.append

        manager.peek_slide_image(0)

        # 다른 파일 변환 중이라도 요청을 버리지 않고 뒤에 붙인다
        assert len(queued) == 1
        assert manager._pending_conversions == 2

        manager.peek_slide_image(1)  # 같은 파일 — 중복 예약 안 함
        assert len(queued) == 1

    def test_counter_returns_to_zero_after_finish(self, manager):
        manager._pending_conversions = 1
        manager._loading = True
        manager._queued_conversions = {Path("/a.pptx")}

        manager._on_single_load_finished(3)

        assert manager._loading is False
        assert manager._queued_conversions == set()
